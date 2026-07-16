from __future__ import annotations

import base64
import hashlib
import json
import os
import re
from io import BytesIO
from pathlib import Path
from typing import BinaryIO

try:
    from models import Evidence, EvidenceModality
except ImportError:
    # 如果无法导入，定义简单的替代类
    class Evidence:
        def __init__(self, **kwargs):
            for k, v in kwargs.items():
                setattr(self, k, v)

    class StrEnumMock:
        """模拟 Enum 成员，提供 .value 属性以兼容真实 EvidenceModality 的调用方式"""
        def __init__(self, value: str):
            self.value = value

        def __str__(self) -> str:
            return self.value

        def __repr__(self) -> str:
            return f"StrEnumMock({self.value!r})"

    class EvidenceModality:
        TEXT = StrEnumMock('text')
        CODE = StrEnumMock('code')
        IMAGE = StrEnumMock('image')


# ---------------------------------------------------------------------------
# 多模态视觉文档 RAG 管道 (Task 1)
# ---------------------------------------------------------------------------

_VISION_LLM_AVAILABLE: bool | None = None  # 延迟检测缓存


def _check_vision_llm() -> bool:
    """检测是否配置了多模态视觉大模型，结果缓存到模块级变量。"""
    global _VISION_LLM_AVAILABLE
    if _VISION_LLM_AVAILABLE is not None:
        return _VISION_LLM_AVAILABLE
    try:
        from config import CONFIG
        if CONFIG.multimodal_llm_api_key and CONFIG.multimodal_llm_endpoint:
            _VISION_LLM_AVAILABLE = True
            return True
        if os.getenv("ZHIPUAI_API_KEY", ""):
            _VISION_LLM_AVAILABLE = True
            return True
    except Exception:
        pass
    _VISION_LLM_AVAILABLE = False
    return False


def _render_pdf_to_images(raw: bytes, dpi: int = 150) -> list[dict]:
    """使用 PyMuPDF 将 PDF 每一页渲染为 PNG 图像。

    Args:
        raw: PDF 文件原始字节
        dpi: 渲染分辨率

    Returns:
        [{page: int, image_bytes: bytes, width: int, height: int}, ...]
        如果 PyMuPDF 不可用或渲染失败，返回空列表
    """
    pages = []
    try:
        import fitz  # PyMuPDF
        doc = fitz.open(stream=raw, filetype="pdf")
        zoom = dpi / 72.0  # PDF 默认 72 DPI
        matrix = fitz.Matrix(zoom, zoom)
        for page_num in range(len(doc)):
            page = doc.load_page(page_num)
            pix = page.get_pixmap(matrix=matrix)
            pages.append({
                "page": page_num + 1,
                "image_bytes": pix.tobytes("png"),
                "width": pix.width,
                "height": pix.height,
            })
        doc.close()
    except ImportError:
        pass  # fitz 未安装
    except Exception as e:
        print(f"  [document_parser] PDF 页面渲染失败: {e}")
    return pages


def _describe_image_with_multimodal_llm(image_base64: str, page_num: int, filename: str) -> str | None:
    """调用多模态视觉大模型描述 PDF 页面图片内容。

    Args:
        image_base64: Base64 编码的图片
        page_num: 页码（用于提示词）
        filename: 文件名（用于提示词）

    Returns:
        大模型生成的图片语义描述；如果调用失败返回 None
    """
    if not _check_vision_llm():
        return None
    try:
        import requests
        from config import CONFIG

        endpoint = CONFIG.multimodal_llm_endpoint
        api_key = CONFIG.multimodal_llm_api_key
        model = CONFIG.multimodal_llm_model or "glm-4v"

        if not endpoint or not api_key:
            # 尝试智谱 GLM-4v 兜底
            zhipu_key = os.getenv("ZHIPUAI_API_KEY", "")
            if zhipu_key:
                endpoint = "https://open.bigmodel.cn/api/paas/v4/chat/completions"
                api_key = zhipu_key
                model = "glm-4v"
            else:
                return None

        prompt = (
            f"你是一位知识图谱助教。请仔细查看这份教材 PDF ({filename}) 的第 {page_num} 页，"
            f"用 2-4 句话描述该页面的核心教学内容，包括：\n"
            f"1. 页面主要讲解的知识点/概念\n"
            f"2. 出现的公式、图表或代码的作用\n"
            f"3. 与机器学习课程中哪些前置知识相关\n\n"
            f"请用中文直接回答，不要加前缀。"
        )

        payload = {
            "model": model,
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {"type": "text", "text": prompt},
                        {"type": "image_url", "image_url": {"url": f"data:image/png;base64,{image_base64}"}},
                    ],
                }
            ],
            "temperature": 0.2,
            "max_tokens": 1024,
        }
        headers = {
            "Content-Type": "application/json",
            "Authorization": f"Bearer {api_key}",
        }
        resp = requests.post(endpoint, json=payload, headers=headers, timeout=60)
        resp.raise_for_status()
        data = resp.json()
        return data["choices"][0]["message"]["content"].strip()
    except ImportError:
        return None  # requests 未安装
    except Exception as e:
        print(f"  [document_parser] 多模态大模型图片描述失败 (第{page_num}页): {e}")
        return None


def _describe_image_with_pil(image_bytes: bytes) -> str:
    """使用 PIL 生成基础图片元数据描述（非视觉模型时的降级方案）。"""
    try:
        from PIL import Image
        img = Image.open(BytesIO(image_bytes))
        w, h = img.size
        mode = img.mode
        fmt = img.format or "PNG"
        dominant = _get_dominant_colors(img)
        return f"教材图片 {fmt} {w}x{h} 模式:{mode} 主色调:{dominant}"
    except Exception:
        return f"教材图片 {len(image_bytes)} bytes"


def parse_pdf_visually(raw: bytes, filename: str) -> list[Evidence]:
    """解析 PDF 的视觉内容：逐页渲染为图片 → 调用多模态大模型生成语义描述。

    Args:
        raw: PDF 文件原始字节
        filename: 文件名（用于生成 Evidence ID）

    Returns:
        Evidence 列表（modality=IMAGE），包含视觉描述和页面图片路径
    """
    pages = _render_pdf_to_images(raw, dpi=150)
    if not pages:
        return []

    # 确保 patches 目录存在
    patches_dir = Path("data/patches")
    patches_dir.mkdir(parents=True, exist_ok=True)

    evidence_list = []
    base_name = Path(filename).stem.replace("..", "").replace("/", "").replace(chr(92), "")[:40]
    vision_available = _check_vision_llm()

    for page_info in pages:
        page_num = page_info["page"]
        img_bytes = page_info["image_bytes"]

        # 保存页面图片到 patches 目录
        img_filename = f"{base_name}_p{page_num}.png"
        img_path = patches_dir / img_filename
        with open(img_path, "wb") as f:
            f.write(img_bytes)

        # 生成图片语义描述
        if vision_available:
            img_b64 = base64.b64encode(img_bytes).decode("utf-8")
            description = _describe_image_with_multimodal_llm(img_b64, page_num, filename)
            if not description:
                description = _describe_image_with_pil(img_bytes)
        else:
            description = _describe_image_with_pil(img_bytes)

        # 自动推断标签
        tags = _infer_tags(description)

        evidence_id = hashlib.sha256(f"{filename}:page{page_num}".encode()).hexdigest()[:16]
        evidence_list.append(Evidence(
            id=f"VIS_{evidence_id}",
            title=f"{filename} 第{page_num}页 视觉解析",
            content=description,
            modality=EvidenceModality.IMAGE,
            source=f"PDF视觉解析:{filename}",
            tags=tags,
            anchors=tuple(),
            metadata={
                "raw_image_ref": str(img_path),
                "page": page_num,
                "width": page_info["width"],
                "height": page_info["height"],
                "doc_source": filename,
            },
        ))

    return evidence_list


# ---------------------------------------------------------------------------
# 原有文档解析逻辑
# ---------------------------------------------------------------------------

def parse_uploaded_file(file: BinaryIO, filename: str) -> str:
    ext = Path(filename).suffix.lower()
    raw = file.read()
    if ext == ".pdf":
        return _parse_pdf(raw)
    elif ext == ".docx":
        return _parse_docx(raw)
    elif ext == ".pptx":
        return _parse_pptx(raw)
    elif ext == ".md":
        return raw.decode("utf-8", errors="replace")
    elif ext == ".txt":
        return raw.decode("utf-8", errors="replace")
    elif ext == ".py":
        try:
            # Python 代码本质上就是文本，直接用 utf-8 读取
            content = raw.decode("utf-8", errors="replace")
            # 在开头加个标记，告诉大模型这是 Python 代码
            return f"【这是 Python 代码文件】\n{content}"
        except Exception as e:
            print(f"读取代码文件失败: {e}")
            return ""
    elif ext in (".mp4", ".avi", ".mov", ".mkv", ".webm", ".flv"):
        return _transcribe_video(raw, filename)
    else:
        return raw.decode("utf-8", errors="replace")


def _parse_pdf(raw: bytes) -> str:
    """解析 PDF 文本，优先 pdfplumber（支持表格提取），降级 PyPDF2，保底 decode。
    
    返回包含 Markdown 表格和 LaTeX 公式的纯净文本。
    """
    try:
        import pdfplumber
        pages_text = []
        with pdfplumber.open(BytesIO(raw)) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                page_parts = [f"--- Page {page_num} ---"]
                
                # 提取文本（保留段落间距）
                text = page.extract_text(x_tolerance=3, y_tolerance=3) or ""
                if text.strip():
                    page_parts.append(text.strip())
                
                # 提取表格并转为 Markdown 格式
                tables = page.extract_tables()
                for table in tables:
                    if not table:
                        continue
                    md_lines = []
                    # 表头
                    header = [_markdown_escape_cell(c) for c in (table[0] or [])]
                    if header:
                        md_lines.append("| " + " | ".join(header) + " |")
                        md_lines.append("| " + " | ".join("---" for _ in header) + " |")
                    # 数据行
                    for row in table[1:]:
                        cells = [_markdown_escape_cell(c) for c in (row or [])]
                        if cells:
                            md_lines.append("| " + " | ".join(cells) + " |")
                    if md_lines:
                        page_parts.append("\n".join(md_lines))
                
                pages_text.append("\n\n".join(page_parts))
        return "\n\n".join(pages_text)
    except Exception:
        pass
    try:
        import PyPDF2
        reader = PyPDF2.PdfReader(BytesIO(raw))
        return "\n".join(page.extract_text() or "" for page in reader.pages)
    except Exception:
        pass
    return raw.decode("utf-8", errors="replace")


def _markdown_escape_cell(value: str | None) -> str:
    """转义表格单元格中可能破坏 Markdown 表格的字符。"""
    if value is None:
        return ""
    return value.replace("|", "\\|").replace("\n", " ").strip()


def _parse_docx(raw: bytes) -> str:
    try:
        import docx
        doc = docx.Document(BytesIO(raw))
        full_text = []
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                full_text.append(text)
        for table in doc.tables:
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                full_text.append(" | ".join(cells))
        return "\n".join(full_text)
    except Exception:
        pass
    return raw.decode("utf-8", errors="replace")


def _parse_pptx(raw: bytes) -> str:
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation(BytesIO(raw))
        pages = []
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_texts = [f"--- 幻灯片 {slide_num} ---"]
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        slide_texts.append(" | ".join(cells))
                if hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        img_desc = _describe_image(img_bytes)
                        if img_desc:
                            slide_texts.append(f"[图片描述] {img_desc}")
                    except Exception:
                        pass
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_texts.append(f"[演讲者备注] {notes_text}")
            pages.append("\n".join(slide_texts))
        return "\n\n".join(pages)
    except Exception:
        pass
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(BytesIO(raw)) as z:
            texts = []
            for name in z.namelist():
                if name.startswith("ppt/slides/slide") and name.endswith(".xml"):
                    content = z.read(name)
                    root = ET.fromstring(content)
                    ns = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}
                    slide_texts = []
                    for t in root.iter("{http://schemas.openxmlformats.org/drawingml/2006/main}t"):
                        if t.text:
                            slide_texts.append(t.text)
                    if slide_texts:
                        texts.append("--- 幻灯片 ---\n" + "\n".join(slide_texts))
            return "\n\n".join(texts) if texts else raw.decode("utf-8", errors="replace")
    except Exception:
        pass
    return raw.decode("utf-8", errors="replace")


def _describe_image(img_bytes: bytes) -> str:
    try:
        import io
        from PIL import Image
        img = Image.open(io.BytesIO(img_bytes))
        w, h = img.size
        mode = img.mode
        format_name = img.format or "unknown"
        dominant_colors = _get_dominant_colors(img)
        return f"图片 {format_name} {w}x{h} 模式:{mode} 主色调:{dominant_colors}"
    except ImportError:
        return f"图片 {len(img_bytes)} bytes (需安装Pillow以解析)"


def _get_dominant_colors(img, n=3):
    try:
        reduced = img.resize((32, 32))
        palette = reduced.getcolors(1024)
        if not palette:
            return "未知"
        palette.sort(key=lambda x: x[0], reverse=True)
        colors = []
        for count, color in palette[:n]:
            colors.append(f"rgb{color[:3]}")
        return ", ".join(colors)
    except Exception:
        return "未知"


def _transcribe_video(raw: bytes, filename: str) -> str:
    temp_path = None
    try:
        import tempfile
        temp_dir = Path("data/uploads/temp_videos")
        temp_dir.mkdir(parents=True, exist_ok=True)
        temp_path = temp_dir / Path(filename).name
        with open(temp_path, "wb") as f:
            f.write(raw)

        transcription = _audio_transcription(str(temp_path))
        if transcription:
            return transcription

        try:
            from moviepy import VideoFileClip
            audio_path = str(temp_path) + ".wav"
            with VideoFileClip(str(temp_path)) as video:
                if video.audio is None:
                    return f"[视频 {filename}] 无音轨，仅元数据: {video.duration:.1f}s {video.size}"
                video.audio.write_audiofile(audio_path, logger=None)
            with open(audio_path, "rb") as af:
                audio_bytes = af.read()
            os.remove(audio_path)
            text = _speech_to_text(audio_bytes)
            duration = "(语音转文字完成)"
            return f"[视频 {filename}] {duration}\n{text}" if text else f"[视频 {filename}] 语音转文字未返回结果"
        except Exception:
            pass

        return f"[视频 {filename}] 需安装 moviepy + whisper 实现语音转文字"
    except Exception as e:
        return f"[视频 {filename}] 转写失败: {e}"
    finally:
        if temp_path and temp_path.exists():
            try:
                os.remove(temp_path)
            except Exception:
                pass


def _audio_transcription(audio_path: str) -> str:
    try:
        import whisper
        model = whisper.load_model("base")
        result = model.transcribe(audio_path, language="zh")
        segments = result.get("segments", [])
        texts = []
        for seg in segments:
            start = seg.get("start", 0)
            end = seg.get("end", 0)
            text = seg.get("text", "").strip()
            if text:
                texts.append(f"[{start:.1f}s-{end:.1f}s] {text}")
        return "\n".join(texts) if texts else result.get("text", "")
    except ImportError:
        pass
    try:
        import subprocess
        result = subprocess.run(
            ["whisper", audio_path, "--model", "base", "--output_format", "txt"],
            capture_output=True, text=True, timeout=300
        )
        if result.returncode == 0:
            return result.stdout.strip()
    except Exception:
        pass
    return ""


def _speech_to_text(audio_bytes: bytes) -> str:
    try:
        import io
        import whisper
        model = whisper.load_model("base")
        import tempfile
        with tempfile.NamedTemporaryFile(suffix=".wav", delete=False) as f:
            f.write(audio_bytes)
            tmp_path = f.name
        try:
            result = model.transcribe(tmp_path, language="zh")
            return result.get("text", "")
        finally:
            os.unlink(tmp_path)
    except ImportError:
        pass
    return ""


def parse_pptx_slides(raw: bytes) -> tuple[dict, ...]:
    slides = []
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        prs = Presentation(BytesIO(raw))
        for slide_num, slide in enumerate(prs.slides, 1):
            texts = []
            images = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        t = para.text.strip()
                        if t:
                            texts.append(t)
                if hasattr(shape, "image"):
                    try:
                        img_bytes = shape.image.blob
                        import base64
                        img_b64 = base64.b64encode(img_bytes).decode("utf-8")
                        images.append({"data": img_b64, "content_type": shape.image.content_type})
                    except Exception:
                        pass
            slides.append({
                "slide_num": slide_num,
                "text": "\n".join(texts),
                "images": images,
                "notes": slide.notes_slide.notes_text_frame.text.strip() if slide.has_notes_slide else "",
            })
    except Exception:
        pass
    return tuple(slides)


def chunk_document(text: str, source: str, chunk_size: int = 520, overlap: int = 80) -> tuple[Evidence, ...]:
    """将文档切分为父子分块结构。

    对于非代码文档：
      -父块 (Parent): ~1000-1500 字，按段落边界切分
      -子块 (Child): ~200-250 字，在父块内部切分
      -每个子块的 Evidence.metadata 包含 parent_content 供 RAG 替换
    对于代码文件：保持原有类和函数边界切分逻辑。
    """
    if not text.strip():
        return ()
    
    # 判断是否是代码文件
    is_code_file = source.lower().endswith('.py')
    
    if is_code_file:
        # 代码文件：优先按类和函数边界切分（保持原有逻辑）
        return _chunk_code(text, source)
    else:
        # 普通文档：父子分块
        return _chunk_with_parent_child(text, source, child_size=chunk_size)


def _chunk_code(text: str, source: str) -> tuple[Evidence, ...]:
    """代码文件分块（保持原有逻辑）。"""
    separators = ["\nclass ", "\ndef ", "\n\n", "\n", " ", ""]
    chunk_size = 1000
    overlap = 200
    text = text.strip()
    
    chunks: list[Evidence] = []
    start = 0
    idx = 1
    text_len = len(text)
    
    while start < text_len:
        end = min(text_len, start + chunk_size)
        content = text[start:end]
        
        if end < text_len:
            best_split = end
            for sep in separators:
                pos = content.rfind(sep)
                if pos > chunk_size // 2:
                    best_split = start + pos + len(sep)
                    break
            end = best_split
        
        content = text[start:end].strip()
        if not content:
            start += 1
            continue
            
        chunk_id = hashlib.sha256(f"{source}:{idx}:{content[:64]}".encode()).hexdigest()[:16]
        chunks.append(Evidence(
            id=chunk_id,
            title=f"{source} chunk {idx}",
            content=content,
            modality=EvidenceModality.TEXT,
            source=source,
            tags=_infer_tags(content),
            anchors=tuple(),
            metadata={"chunk_index": idx, "doc_source": source},
        ))
        if end == text_len:
            break
        start = max(0, end - overlap)
        idx += 1
    return tuple(chunks)


def _chunk_with_parent_child(text: str, source: str, child_size: int = 250) -> tuple[Evidence, ...]:
    """父子分块：先切父块(~1000-1500字)，再在每个父块内切子块(~200-250字)。"""
    if not text.strip():
        return ()
    
    # 保留原始段落: 按双换行切分出自然段落
    paragraphs = [p.strip() for p in text.split("\n\n") if p.strip()]
    if not paragraphs:
        paragraphs = [text.strip()]
    
    # 合并段落为父块 (~1000-1500 字)
    parents: list[str] = []
    current_parent = ""
    for para in paragraphs:
        if len(current_parent) + len(para) < 1200:
            current_parent = (current_parent + "\n\n" + para).strip()
        else:
            if current_parent:
                parents.append(current_parent)
            current_parent = para
    if current_parent:
        parents.append(current_parent)
    
    # 如果父块太大(>2000字)，按句号/换行再分割
    refined_parents: list[str] = []
    for p in parents:
        if len(p) > 2000:
            # 尝试按句号分割
            sub_parts = re.split(r'(?<=[。！？.!?])\s+', p)
            buffer = ""
            for sp in sub_parts:
                if len(buffer) + len(sp) < 1200:
                    buffer = (buffer + sp).strip()
                else:
                    if buffer:
                        refined_parents.append(buffer)
                    buffer = sp
            if buffer:
                refined_parents.append(buffer)
        else:
            refined_parents.append(p)
    parents = refined_parents
    
    children: list[Evidence] = []
    for parent_idx, parent_text in enumerate(parents, 1):
        # 在父块内切子块 (~200-250 字)
        # 优先在句子边界切割
        sentences = re.split(r'(?<=[。！？.!?\n])\s*', parent_text)
        sentences = [s.strip() for s in sentences if s.strip()]
        
        child_idx = 1
        child_buffer = ""
        for sent in sentences:
            if len(child_buffer) + len(sent) < child_size:
                child_buffer = (child_buffer + " " + sent).strip()
            else:
                if child_buffer:
                    chunk_id = hashlib.sha256(f"{source}:p{parent_idx}c{child_idx}:{child_buffer[:64]}".encode()).hexdigest()[:16]
                    children.append(Evidence(
                        id=chunk_id,
                        title=f"{source} P{parent_idx}C{child_idx}",
                        content=child_buffer,
                        modality=EvidenceModality.TEXT,
                        source=source,
                        tags=_infer_tags(child_buffer),
                        anchors=tuple(),
                        metadata={
                            "chunk_index": child_idx,
                            "parent_index": parent_idx,
                            "child_index": child_idx,
                            "doc_source": source,
                            "parent_content": parent_text,
                        },
                    ))
                    child_idx += 1
                child_buffer = sent
        if child_buffer:
            chunk_id = hashlib.sha256(f"{source}:p{parent_idx}c{child_idx}:{child_buffer[:64]}".encode()).hexdigest()[:16]
            children.append(Evidence(
                id=chunk_id,
                title=f"{source} P{parent_idx}C{child_idx}",
                content=child_buffer,
                modality=EvidenceModality.TEXT,
                source=source,
                tags=_infer_tags(child_buffer),
                anchors=tuple(),
                metadata={
                    "chunk_index": child_idx,
                    "parent_index": parent_idx,
                    "child_index": child_idx,
                    "doc_source": source,
                    "parent_content": parent_text,
                },
            ))
    
    return tuple(children)


def _infer_tags(content: str) -> tuple[str, ...]:
    tags: list[str] = []
    topic_map = {
        "池化": "池化层", "卷积": "卷积", "梯度": "梯度下降",
        "反向传播": "反向传播", "链式法则": "链式法则",
        "逻辑回归": "逻辑回归", "线性回归": "线性回归",
        "过拟合": "过拟合", "正则化": "正则化",
        "交叉验证": "交叉验证", "混淆矩阵": "混淆矩阵",
        "机器学习": "机器学习", "分类": "分类", "回归": "回归",
        "聚类": "聚类", "神经网络": "神经网络",
        "Python": "Python", "numpy": "NumPy", "pandas": "Pandas",
        "matplotlib": "Matplotlib", "scikit-learn": "ScikitLearn",
        "CNN": "CNN", "RNN": "RNN", "transformer": "Transformer",
        "注意力": "注意力机制", "损失函数": "损失函数",
        "激活函数": "激活函数", "归一化": "归一化",
    }
    seen = set()
    for keyword, tag in topic_map.items():
        if keyword in content and tag not in seen:
            tags.append(tag)
            seen.add(tag)
    return tuple(tags[:8])