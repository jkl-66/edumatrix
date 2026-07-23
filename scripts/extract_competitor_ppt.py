from __future__ import annotations

import argparse
import hashlib
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


def iter_shapes(shapes):
    for shape in shapes:
        yield shape
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            yield from iter_shapes(shape.shapes)


def extract_text(shape) -> list[str]:
    texts: list[str] = []
    if getattr(shape, "has_text_frame", False):
        for paragraph in shape.text_frame.paragraphs:
            text = "".join(run.text for run in paragraph.runs).strip()
            if not text:
                text = paragraph.text.strip()
            if text:
                texts.append(text)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            values = [cell.text.strip() for cell in row.cells]
            if any(values):
                texts.append(" | ".join(values))
    return texts


def notes_text(slide) -> str:
    try:
        frame = slide.notes_slide.notes_text_frame
        return frame.text.strip() if frame else ""
    except Exception:
        return ""


def save_picture(shape, media_dir: Path, slide_no: int, shape_no: int) -> dict[str, Any] | None:
    if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
        return None
    image = shape.image
    digest = hashlib.sha256(image.blob).hexdigest()
    target = media_dir / f"slide-{slide_no:02d}-shape-{shape_no:02d}-{digest[:10]}.{image.ext}"
    if not target.exists():
        target.write_bytes(image.blob)
    return {
        "path": target.as_posix(),
        "sha256": digest,
        "content_type": image.content_type,
        "size_bytes": len(image.blob),
        "width_emu": int(shape.width),
        "height_emu": int(shape.height),
    }


def extract(input_path: Path, output_dir: Path) -> dict[str, Any]:
    output_dir.mkdir(parents=True, exist_ok=True)
    media_dir = output_dir / "media"
    media_dir.mkdir(parents=True, exist_ok=True)

    presentation = Presentation(str(input_path))
    result: dict[str, Any] = {
        "source": str(input_path.resolve()),
        "slide_count": len(presentation.slides),
        "slide_width_emu": int(presentation.slide_width),
        "slide_height_emu": int(presentation.slide_height),
        "slides": [],
    }

    for slide_no, slide in enumerate(presentation.slides, start=1):
        slide_data: dict[str, Any] = {
            "slide": slide_no,
            "texts": [],
            "tables": [],
            "pictures": [],
            "notes": notes_text(slide),
        }
        for shape_no, shape in enumerate(iter_shapes(slide.shapes), start=1):
            texts = extract_text(shape)
            if texts:
                slide_data["texts"].extend(texts)
            if getattr(shape, "has_table", False):
                slide_data["tables"].append(
                    [[cell.text.strip() for cell in row.cells] for row in shape.table.rows]
                )
            picture = save_picture(shape, media_dir, slide_no, shape_no)
            if picture:
                slide_data["pictures"].append(picture)
        result["slides"].append(slide_data)

    (output_dir / "slides.json").write_text(
        json.dumps(result, ensure_ascii=False, indent=2), encoding="utf-8"
    )

    lines = [
        "# 57016115 介绍 PPT 逐页证据",
        "",
        f"- 来源：`{input_path}`",
        f"- 页数：{result['slide_count']}",
        "- 说明：本文件由 PPTX 结构化解析生成，文字顺序按形状遍历顺序记录。",
        "",
    ]
    for slide in result["slides"]:
        lines.extend([f"## 第 {slide['slide']} 页", ""])
        if slide["texts"]:
            lines.extend(f"- {text}" for text in slide["texts"])
        else:
            lines.append("- （未提取到文字）")
        if slide["pictures"]:
            lines.extend(["", "图片证据：", ""])
            lines.extend(f"- `{picture['path']}`" for picture in slide["pictures"])
        if slide["notes"]:
            lines.extend(["", "备注：", "", slide["notes"]])
        lines.append("")
    (output_dir / "slides.md").write_text("\n".join(lines), encoding="utf-8")
    return result


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("input", type=Path)
    parser.add_argument("output", type=Path)
    args = parser.parse_args()
    result = extract(args.input, args.output)
    print(f"slides={result['slide_count']}")
    print(args.output.resolve())


if __name__ == "__main__":
    main()
