"""Build a compact competition PPT from verified local evidence."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "outputs" / "EduMatrix_软件杯答辩PPT_待补团队信息.pptx"
NOTES = ROOT / "outputs" / "EduMatrix_软件杯PPT逐页讲解备注.md"
SCREENSHOTS = ROOT / "outputs" / "e2e_no_docker"

NAVY = RGBColor(25, 37, 34)
INK = RGBColor(30, 40, 34)
MUTED = RGBColor(99, 111, 103)
GREEN = RGBColor(48, 75, 61)
MINT = RGBColor(222, 233, 225)
BLUE = RGBColor(75, 104, 218)
AMBER = RGBColor(225, 154, 54)
WHITE = RGBColor(255, 255, 255)
LIGHT = RGBColor(245, 247, 244)


def add_box(slide, x, y, w, h, fill=WHITE, line=None, radius=False):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    if radius:
        shape.adjustments[0] = 0.12
    return shape


def add_text(slide, text, x, y, w, h, size=18, color=INK, bold=False, align=PP_ALIGN.LEFT, font="Microsoft YaHei"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.06)
    frame.margin_right = Inches(0.06)
    frame.margin_top = Inches(0.03)
    frame.margin_bottom = Inches(0.03)
    p = frame.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_bullets(slide, items, x, y, w, h, size=15, color=INK, bullet_color=BLUE):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.04)
    for index, item in enumerate(items):
        p = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        p.text = f"●  {item}"
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = color
    return box


def add_header(slide, kicker, title, number):
    add_text(slide, kicker.upper(), 0.55, 0.28, 5.0, 0.28, size=9, color=BLUE, bold=True)
    add_text(slide, title, 0.55, 0.60, 11.7, 0.55, size=25, color=NAVY, bold=True)
    add_text(slide, f"{number:02d}", 12.25, 0.35, 0.55, 0.4, size=15, color=MUTED, bold=True, align=PP_ALIGN.RIGHT)
    add_box(slide, 0.55, 1.25, 12.2, 0.015, fill=MINT)


def add_footer(slide):
    add_text(slide, "EduMatrix 智教矩阵  |  软件杯 A3  |  事实状态以源码与证据为准", 0.55, 7.18, 11.8, 0.2, size=8, color=MUTED)


def add_card(slide, title, body, x, y, w, h, accent=BLUE):
    add_box(slide, x, y, w, h, fill=WHITE, line=RGBColor(224, 230, 225), radius=True)
    add_box(slide, x, y, 0.07, h, fill=accent, line=accent, radius=True)
    add_text(slide, title, x + 0.22, y + 0.18, w - 0.35, 0.35, size=15, color=NAVY, bold=True)
    add_text(slide, body, x + 0.22, y + 0.63, w - 0.38, h - 0.78, size=11, color=MUTED)


def add_image_if_exists(slide, filename, x, y, w, h):
    path = SCREENSHOTS / filename
    if path.exists():
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
    else:
        add_box(slide, x, y, w, h, fill=MINT, line=MINT, radius=True)
        add_text(slide, f"待补截图\n{filename}", x, y + h / 2 - 0.2, w, 0.5, size=12, color=MUTED, align=PP_ALIGN.CENTER)


def new_slide(prs, title=None, kicker="EduMatrix"):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = LIGHT
    if title:
        add_header(slide, kicker, title, len(prs.slides))
    add_footer(slide)
    return slide


def build() -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1 Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = NAVY
    add_text(slide, "软件杯 A3 · 作品答辩", 0.7, 0.7, 4.0, 0.3, size=13, color=MINT, bold=True)
    add_text(slide, "EduMatrix\n智教矩阵", 0.7, 1.45, 5.7, 1.45, size=38, color=WHITE, bold=True)
    add_text(slide, "面向机器学习导论的个性化资源生成与学习多智能体系统", 0.75, 3.15, 5.5, 0.65, size=17, color=MINT)
    add_box(slide, 0.75, 4.35, 3.15, 0.52, fill=RGBColor(75, 112, 87), line=RGBColor(75, 112, 87), radius=True)
    add_text(slide, "一问 · 多智能体 · 五类资源 · 反馈闭环", 0.9, 4.49, 2.85, 0.2, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_image_if_exists(slide, "02-dashboard.png", 7.0, 0.85, 5.7, 5.95)
    add_text(slide, "团队 / 学校：待补\n版本：2026-07-19", 0.75, 6.45, 4.0, 0.45, size=10, color=RGBColor(184, 201, 190))

    # 2 Problem
    slide = new_slide(prs, "为什么需要 EduMatrix？", "问题定义")
    add_text(slide, "同一份课程资源，无法同时适配不同基础、目标和认知偏好的学习者。", 0.7, 1.65, 11.7, 0.5, size=20, color=NAVY, bold=True)
    add_card(slide, "资源同质化", "讲义、代码和练习通常由学生自己拼接，难以根据学习状态改变顺序和难度。", 0.7, 2.55, 3.8, 2.15, BLUE)
    add_card(slide, "知识缺证据", "大模型回答缺少课程图谱、来源和跨模态一致性约束，理科内容更容易产生认知冲突。", 4.78, 2.55, 3.8, 2.15, AMBER)
    add_card(slide, "反馈不闭环", "答题、错因、掌握度和复习计划分散在不同页面，下一次学习无法使用最新状态。", 8.86, 2.55, 3.8, 2.15, GREEN)
    add_text(slide, "目标：把“学生提问”升级为可观测、可解释、可持续更新的学习动作。", 0.75, 5.45, 11.6, 0.5, size=18, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 3 Loop
    slide = new_slide(prs, "一条完整的个性化学习闭环", "产品闭环")
    stages = [("01", "画像诊断", "专业、目标、偏好、掌握度"), ("02", "路径规划", "知识图谱 + ZPD + A*"), ("03", "证据检索", "文本、图像、图谱来源"), ("04", "资源生成", "讲义、导图、代码、题目、视频"), ("05", "反馈更新", "答题、错因、复习与再规划")]
    x = 0.65
    for i, (num, title, body) in enumerate(stages):
        add_box(slide, x, 2.2, 2.25, 2.0, fill=WHITE, line=RGBColor(222, 230, 224), radius=True)
        add_text(slide, num, x + 0.17, 2.42, 0.5, 0.35, size=17, color=BLUE, bold=True)
        add_text(slide, title, x + 0.17, 2.95, 1.9, 0.35, size=16, color=NAVY, bold=True)
        add_text(slide, body, x + 0.17, 3.48, 1.85, 0.55, size=11, color=MUTED)
        if i < 4:
            add_text(slide, "→", x + 2.29, 2.92, 0.5, 0.35, size=24, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
        x += 2.55
    add_box(slide, 1.45, 5.1, 10.4, 0.7, fill=MINT, line=MINT, radius=True)
    add_text(slide, "每一次交互都会沉淀为下一次决策的上下文，而不是一次性回答。", 1.7, 5.29, 9.9, 0.3, size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 4 Architecture
    slide = new_slide(prs, "系统架构：1 + 3 + 5 协同矩阵", "技术架构")
    add_card(slide, "交互中枢 · 1", "苏格拉底导师\n接收问题、统一路由、流式输出", 0.75, 1.85, 3.6, 1.35, BLUE)
    add_card(slide, "认知治理 · 3", "画像探针 · 路径规划师 · 量化评估师\n把学习状态转成决策约束", 4.85, 1.85, 3.6, 1.35, GREEN)
    add_card(slide, "资源工厂 · 5", "理论教授 · 逻辑画师 · 极客助教\n考官智能体 · 视频推荐官", 8.95, 1.85, 3.6, 1.35, AMBER)
    add_box(slide, 1.1, 4.0, 11.1, 1.25, fill=NAVY, line=NAVY, radius=True)
    add_text(slide, "Vue 3 / Vite", 1.55, 4.42, 1.7, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↔", 3.25, 4.39, 0.55, 0.35, size=21, color=MINT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "FastAPI + JWT + SSE", 3.8, 4.42, 2.3, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↔", 6.1, 4.39, 0.55, 0.35, size=21, color=MINT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "Hybrid RAG + 学情算法", 6.7, 4.42, 2.45, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "↔", 9.2, 4.39, 0.55, 0.35, size=21, color=MINT, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "SQLite/WAL + 可选沙箱", 9.75, 4.42, 2.1, 0.3, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, "设计原则：证据进入生成前，跨模态一致性检查发生在生成后。", 1.1, 5.85, 11.1, 0.35, size=16, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 5 Personalization
    slide = new_slide(prs, "同一主题，三种学习路径", "个性化")
    add_image_if_exists(slide, "06-learning-path.png", 0.7, 1.65, 7.0, 4.95)
    add_card(slide, "视觉型学习者", "图示演示 + 具体例子\n优先建立概念关系和可视化路径", 8.0, 1.75, 4.5, 1.2, BLUE)
    add_card(slide, "项目型学习者", "代码实操 + 分步引导\n优先进入可运行案例和迁移练习", 8.0, 3.18, 4.5, 1.2, GREEN)
    add_card(slide, "研究型学习者", "对比解析 + 分步引导\n优先解释前置关系和理论边界", 8.0, 4.61, 4.5, 1.2, AMBER)

    # 6 Agent output
    slide = new_slide(prs, "一次提问，生成五类可用资源", "多智能体")
    add_image_if_exists(slide, "04-chat-response.png", 0.7, 1.65, 6.45, 4.95)
    resources = [("讲义", "概念、推导、误区", BLUE), ("导图", "前置关系、结构", GREEN), ("代码", "可操作实训案例", AMBER), ("题目", "分层练习与反馈", BLUE), ("视频", "推荐与脚本", GREEN)]
    y = 1.72
    for title, body, color in resources:
        add_card(slide, title, body, 7.65, y, 4.8, 0.82, color)
        y += 0.98
    add_text(slide, "Agent Timeline 让评委看到协同过程，而不是只看到最终文本。", 7.7, 6.62, 4.6, 0.35, size=12, color=GREEN, bold=True)

    # 7 RAG
    slide = new_slide(prs, "证据链路：让生成结果可追溯", "RAG 与一致性")
    add_card(slide, "检索", "课程知识图谱\n文本切片 / 图像切片\n固定证据 ID 与分数", 0.75, 1.8, 3.6, 2.25, BLUE)
    add_text(slide, "→", 4.43, 2.55, 0.55, 0.45, size=25, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, "清洗", "DRAG 证据裁决\n低置信度拒答\nowner 过滤私有资料", 4.85, 1.8, 3.6, 2.25, AMBER)
    add_text(slide, "→", 8.53, 2.55, 0.55, 0.45, size=25, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, "校验", "公式、变量、概念\n跨模态对齐\n失败时局部重生成", 8.95, 1.8, 3.6, 2.25, GREEN)
    add_box(slide, 1.05, 4.8, 11.15, 0.95, fill=NAVY, line=NAVY, radius=True)
    add_text(slide, "本地固定知识索引实测：返回 6 条证据，含文本与图像模态；外部搜索在证据生成阶段关闭。", 1.35, 5.1, 10.55, 0.35, size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    # 8 Feedback
    slide = new_slide(prs, "反馈不是终点，而是下一轮决策的输入", "学习反馈")
    add_image_if_exists(slide, "06-learning-path.png", 0.75, 1.65, 6.1, 4.95)
    add_bullets(slide, [
        "答题结果 → 掌握度与错因",
        "行为日志 → 认知负荷与专注度",
        "错题记录 → 相似题与间隔复习",
        "新状态 → 路径重规划与资源重排",
    ], 7.25, 2.05, 5.0, 2.6, size=17, color=INK)
    add_box(slide, 7.25, 5.15, 5.0, 0.75, fill=MINT, line=MINT, radius=True)
    add_text(slide, "从“会不会”走向“为什么不会、下一步学什么”。", 7.48, 5.37, 4.55, 0.3, size=15, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    # 9 Deployment / security
    slide = new_slide(prs, "提交策略：无 Docker 默认，Docker 可选增强", "工程与安全")
    add_card(slide, "评委默认路径", "Python + Node + SQLite\n核心学习闭环可启动\n无需 Docker Desktop", 0.8, 1.8, 5.55, 2.05, GREEN)
    add_card(slide, "可选代码能力", "EDUMATRIX_SANDBOX_MODE=docker\n非特权容器、无网络、资源限制\n未启用时明确 503，不回退宿主进程", 6.95, 1.8, 5.55, 2.05, AMBER)
    add_bullets(slide, [
        "JWT 服务端身份绑定与跨用户 owner 过滤",
        "生产环境拒绝固定 JWT secret",
        "上传和远程下载有大小边界",
        "Playwright 只在 PDF 导出时需要，与 Docker 独立",
    ], 1.0, 4.55, 11.2, 1.6, size=15)

    # 10 Evidence
    slide = new_slide(prs, "当前可复核证据", "验证结果")
    metrics = [("80/80", "完整集成测试", GREEN), ("62/62", "成员专项测试", BLUE), ("47/47", "运行时安全矩阵", AMBER), ("10/10", "安全契约", GREEN), ("5/5", "资源类型覆盖", BLUE), ("6", "本地 RAG 证据", AMBER)]
    x = 0.75
    for value, label, color in metrics:
        add_box(slide, x, 1.85, 1.85, 1.4, fill=WHITE, line=RGBColor(224, 230, 225), radius=True)
        add_text(slide, value, x, 2.1, 1.85, 0.45, size=24, color=color, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, label, x + 0.1, 2.75, 1.65, 0.25, size=10, color=MUTED, align=PP_ALIGN.CENTER)
        x += 2.05
    add_text(slide, "证据性质：项目实测 / 代码静态证据 / 合成演示数据。", 0.85, 4.05, 11.5, 0.4, size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_bullets(slide, [
        "无 Docker 浏览器 E2E：注册、初始化、仪表盘、对话、学习路径、沙箱状态",
        "创新证据：三组画像、固定知识集、单 Agent/多 Agent与无 RAG/有 RAG结构对比",
        "未宣称真实用户实验指标，避免把设计目标写成实测结果",
    ], 1.05, 4.75, 11.2, 1.45, size=14)

    # 11 Demo
    slide = new_slide(prs, "七分钟演示路径", "答辩流程")
    timeline = [("0:00", "定位", "问题与价值"), ("0:40", "画像", "三组差异"), ("1:30", "答疑", "同题不同资源"), ("3:00", "证据", "RAG + Timeline"), ("4:20", "反馈", "测验 → 画像 → 路径"), ("6:00", "工程", "安全与部署"), ("6:40", "收束", "证据与边界")]
    y = 1.65
    for time_label, title, body in timeline:
        add_box(slide, 1.0, y, 1.2, 0.48, fill=NAVY, line=NAVY, radius=True)
        add_text(slide, time_label, 1.0, y + 0.12, 1.2, 0.2, size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, title, 2.55, y + 0.03, 1.5, 0.25, size=15, color=NAVY, bold=True)
        add_text(slide, body, 4.15, y + 0.03, 7.9, 0.25, size=14, color=MUTED)
        y += 0.72

    # 12 Closing
    slide = new_slide(prs, "我们交付的不是一个聊天框", "收束")
    add_text(slide, "而是一条可观测、可解释、可复现的学习闭环。", 0.9, 1.85, 11.5, 0.65, size=28, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_card(slide, "现在已具备", "多智能体资源工厂\n个性化画像与路径\n证据与反馈链路\n无 Docker 核心验收", 1.0, 3.15, 3.45, 2.0, GREEN)
    add_card(slide, "正在补齐", "前端最终视觉收口\n官方格式与团队信息\n最终录屏与提交压缩包", 4.95, 3.15, 3.45, 2.0, BLUE)
    add_card(slide, "原则", "数据有标签\n结果可复现\n能力不夸大\n风险主动说明", 8.9, 3.15, 3.45, 2.0, AMBER)
    add_text(slide, "谢谢各位评委", 0.9, 6.1, 11.5, 0.45, size=22, color=GREEN, bold=True, align=PP_ALIGN.CENTER)

    prs.save(OUT)

    notes = """# EduMatrix 软件杯 PPT 逐页讲解备注

> 文件：`EduMatrix_软件杯答辩PPT_待补团队信息.pptx`
> 
> PPT 中的数字仅引用当前工作区已经生成的代码实测或合成演示证据。团队信息、官方模板和最终视频需提交前补齐。

1. **封面**：一句话说明 EduMatrix 面向机器学习导论，核心是个性化资源生成和学习反馈闭环。
2. **问题定义**：强调资源同质化、知识缺证据、反馈不闭环，避免泛泛讲“AI 很重要”。
3. **产品闭环**：按画像、路径、证据、资源、反馈五步讲清楚系统为什么不是单轮问答。
4. **技术架构**：解释 1+3+5 的职责分工，指出前端、API、RAG、学情算法和持久化的边界。
5. **个性化**：展示同一主题下三组画像的输入字段差异；不要把合成画像说成真实调研。
6. **多智能体**：展示一次请求产生五类资源和 Agent Timeline，说明并行生产与结果可观察性。
7. **RAG**：展示证据 ID、来源、模态和对齐；明确证据链路不等于已经证明幻觉率达标。
8. **反馈**：用“答题→掌握度→路径”的关系说明系统如何持续更新。
9. **工程安全**：主动回答为什么不强制 Docker：核心闭环降低评委安装成本，代码沙箱是可选增强且不回退宿主进程。
10. **验证结果**：先说测试和安全矩阵，再说证据性质；不要说“系统不存在风险”。
11. **演示流程**：严格控制 7 分钟，外部 LLM 失败时切换 deterministic 结果，Docker 不可用时跳过代码实时运行。
12. **收束**：用“可观测、可解释、可复现”总结工程价值，并主动说明待补事项。
"""
    NOTES.write_text(notes, encoding="utf-8")
    print(f"created: {OUT}")
    print(f"notes: {NOTES}")


if __name__ == "__main__":
    build()
