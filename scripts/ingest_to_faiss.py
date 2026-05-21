"""
EduMatrix FAISS Ingestion Pipeline
Extracts text from all downloaded files and indexes them into FAISS by category.
Categories: courseware(ipynb/md/html), ppt, code, textbook(pdf/md/tex), dataset, image_desc
"""
import hashlib
import json
import re
import sys
import traceback
from pathlib import Path
from collections import defaultdict

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

import nbformat

from models import Evidence, EvidenceModality
from vector_store_faiss import FaissVectorIndex

BASE = Path(r"d:\code\edumatrix")
DATA_DIR = BASE / "data"
FAISS_DIR = BASE / "data" / "faiss_indexes"
FAISS_DIR.mkdir(parents=True, exist_ok=True)

CATEGORIES = {
    "courseware": "课件（Notebook/课程页面/Markdown课件）",
    "ppt": "PPT幻灯片课件",
    "code": "Python代码示例",
    "textbook": "教材/PDF/TeX教材",
    "dataset": "数据集CSV",
    "reference": "参考文档",
}

TOTAL_STATS = defaultdict(lambda: {"files": 0, "chunks": 0})

# ============================================================
# Text extractors for each file type
# ============================================================

def extract_ipynb(path: Path) -> list[dict]:
    """Extract markdown and code cells from Jupyter notebook."""
    chunks = []
    try:
        with open(path, "r", encoding="utf-8") as f:
            nb = nbformat.read(f, as_version=4)
    except Exception:
        return chunks

    md_cells = []
    code_cells = []
    for cell in nb.cells:
        source = cell.source.strip()
        if not source:
            continue
        if cell.cell_type == "markdown":
            md_cells.append(source)
        elif cell.cell_type == "code":
            code_cells.append(source)

    if md_cells:
        combined = "\n".join(md_cells)
        for i, chunk in enumerate(chunk_text(combined, 800, 100)):
            chunks.append({
                "title": f"{path.stem} 课件-{i + 1}",
                "content": chunk,
                "modality": EvidenceModality.TEXT,
                "tags": infer_tags(chunk) + ("notebook", "markdown"),
                "anchors": infer_anchors(chunk),
            })
    if code_cells:
        combined = "\n".join(code_cells)
        for i, chunk in enumerate(chunk_text(combined, 600, 50)):
            chunks.append({
                "title": f"{path.stem} 代码-{i + 1}",
                "content": chunk,
                "modality": EvidenceModality.CODE,
                "tags": infer_tags(chunk) + ("notebook", "code"),
                "anchors": infer_anchors(chunk),
            })
    return chunks


def extract_python(path: Path) -> list[dict]:
    chunks = []
    try:
        text = path.read_text(encoding="utf-8")
    except Exception:
        return chunks
    if not text.strip():
        return chunks
    for i, chunk in enumerate(chunk_text(text, 600, 50)):
        chunks.append({
            "title": f"{path.stem} code-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.CODE,
            "tags": infer_tags(chunk) + ("python", "code"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_markdown(path: Path) -> list[dict]:
    chunks = []
    try:
        text = path.read_text(encoding="utf-8")
    except Exception:
        return chunks
    if not text.strip():
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} md-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("markdown",),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_html(path: Path) -> list[dict]:
    chunks = []
    try:
        text = path.read_text(encoding="utf-8")
    except Exception:
        return chunks
    text = re.sub(r"<[^>]+>", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) < 50:
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} html-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("html", "course_page"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_pdf(path: Path) -> list[dict]:
    chunks = []
    try:
        from pypdf import PdfReader
        reader = PdfReader(path)
        text_pages = []
        for page in reader.pages:
            t = page.extract_text()
            if t:
                text_pages.append(t)
        text = "\n".join(text_pages)
    except Exception:
        return chunks
    if not text.strip() or len(text.strip()) < 50:
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} pdf-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("pdf", "textbook"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_pptx(path: Path) -> list[dict]:
    chunks = []
    try:
        from pptx import Presentation
        prs = Presentation(path)
        slides_text = []
        for slide in prs.slides:
            slide_content = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        if para.text.strip():
                            slide_content.append(para.text.strip())
            if slide_content:
                slides_text.append("\n".join(slide_content))
        text = "\n\n".join(slides_text)
    except Exception:
        return chunks
    if not text.strip():
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} slide-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("pptx", "slide", "ppt"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_csv(path: Path) -> list[dict]:
    chunks = []
    try:
        import pandas as pd
        df = pd.read_csv(path, nrows=100)
        info = [
            f"Dataset: {path.stem}",
            f"Shape: {df.shape[0]} rows x {df.shape[1]} columns",
            f"Columns: {', '.join(df.columns.tolist())}",
            f"Sample rows:",
        ]
        for _, row in df.head(5).iterrows():
            info.append(" | ".join(f"{col}={val}" for col, val in row.items()))
        text = "\n".join(info)
    except Exception:
        return chunks
    chunks.append({
        "title": f"{path.stem} dataset",
        "content": text,
        "modality": EvidenceModality.TEXT,
        "tags": ("dataset", "csv", path.stem) + infer_tags(text),
        "anchors": infer_anchors(text),
    })
    return chunks


def extract_tex(path: Path) -> list[dict]:
    chunks = []
    try:
        text = path.read_text(encoding="utf-8")
    except Exception:
        return chunks
    text = re.sub(r"\\(?:label|ref|cite)\{[^}]*\}", "", text)
    text = re.sub(r"\\[a-zA-Z]+", " ", text)
    text = re.sub(r"\{|\}|\\", " ", text)
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) < 50:
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} tex-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("latex", "tex", "textbook"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


def extract_epub(path: Path) -> list[dict]:
    chunks = []
    try:
        import zipfile
        import xml.etree.ElementTree as ET
        with zipfile.ZipFile(path) as z:
            html_files = [n for n in z.namelist() if n.endswith((".xhtml", ".html", ".htm"))]
            text_parts = []
            for hf in html_files[:20]:
                content = z.read(hf).decode("utf-8", errors="ignore")
                text_parts.append(re.sub(r"<[^>]+>", " ", content))
            text = "\n".join(text_parts)
    except Exception:
        return chunks
    text = re.sub(r"\s+", " ", text).strip()
    if len(text) < 50:
        return chunks
    for i, chunk in enumerate(chunk_text(text, 800, 100)):
        chunks.append({
            "title": f"{path.stem} epub-{i + 1}",
            "content": chunk,
            "modality": EvidenceModality.TEXT,
            "tags": infer_tags(chunk) + ("epub", "ebook"),
            "anchors": infer_anchors(chunk),
        })
    return chunks


# ============================================================
# Utilities
# ============================================================

def chunk_text(text: str, chunk_size: int, overlap: int) -> list[str]:
    chunks = []
    start = 0
    while start < len(text):
        end = min(start + chunk_size, len(text))
        chunks.append(text[start:end].strip())
        if end == len(text):
            break
        start = end - overlap
    return [c for c in chunks if len(c) > 20]


def infer_tags(text: str) -> tuple:
    concepts = (
        "机器学习", "监督学习", "逻辑回归", "混淆矩阵", "过拟合",
        "正则化", "交叉验证", "池化层", "最大池化", "平均池化",
        "卷积", "特征图", "反向传播", "梯度下降", "线性回归",
        "决策树", "支持向量机", "朴素贝叶斯", "PCA", "聚类",
        "神经网络", "CNN", "RNN", "Transformer", "激活函数",
        "损失函数", "精度", "召回率", "F1", "ROC", "AUC",
        "Iris", "iris", "MNIST", "梯度", "归一化", "标准化",
    )
    return tuple(c for c in concepts if c.lower() in text.lower())


def infer_anchors(text: str) -> tuple:
    anchors = re.findall(r"\b[A-Za-z][A-Za-z0-9_.]{2,}\b", text)
    return tuple(dict.fromkeys(anchors[:12]))


def stable_id(source: str, idx: int) -> str:
    digest = hashlib.sha1(f"{source}:{idx}".encode()).hexdigest()[:12]
    return f"FAISS_{digest}"


# ============================================================
# Main ingestion
# ============================================================

def build_category_index(cat_name: str, cat_label: str, files_by_type: dict[str, list[Path]]):
    print(f"\n{'=' * 60}")
    print(f"Building {cat_name} ({cat_label})...")
    index = FaissVectorIndex(f"edumatrix-{cat_name}")
    total_chunks = 0
    total_files = 0

    extractors = {
        ".ipynb": extract_ipynb,
        ".py": extract_python,
        ".md": extract_markdown,
        ".rst": extract_markdown,
        ".html": extract_html,
        ".htm": extract_html,
        ".pdf": extract_pdf,
        ".pptx": extract_pptx,
        ".csv": extract_csv,
        ".tex": extract_tex,
        ".epub": extract_epub,
    }

    for ext, paths in files_by_type.items():
        extractor = extractors.get(ext)
        if not extractor:
            continue
        for fpath in paths:
            try:
                chunks = extractor(fpath)
                if not chunks:
                    continue
                evidence_list = []
                for i, ch in enumerate(chunks):
                    eid = stable_id(str(fpath.relative_to(BASE)), i)
                    evidence_list.append(Evidence(
                        id=eid,
                        title=ch["title"],
                        content=ch["content"],
                        modality=ch["modality"],
                        source=str(fpath.relative_to(BASE)),
                        tags=ch["tags"],
                        anchors=ch["anchors"],
                        metadata={"category": cat_name, "file_type": ext},
                    ))
                index.upsert(tuple(evidence_list))
                total_chunks += len(evidence_list)
                total_files += 1
            except Exception as e:
                print(f"    [SKIP] {fpath.name}: {e}")

    TOTAL_STATS[cat_name]["files"] = total_files
    TOTAL_STATS[cat_name]["chunks"] = total_chunks

    save_path = FAISS_DIR / f"{cat_name}_index"
    index.save(save_path)
    print(f"  -> Indexed {total_files} files -> {total_chunks} chunks")
    print(f"  -> Saved to {save_path}")
    print(f"  -> Index size: {index.count()} vectors")

    return index


def classify_files():
    """Classify all files under data/ into categories."""
    files_by_ext = defaultdict(list)
    for fpath in DATA_DIR.rglob("*"):
        if not fpath.is_file():
            continue
        ext = fpath.suffix.lower()
        skip_exts = {".gitkeep", ".gitignore", ".png", ".jpg", ".jpeg", ".gif", ".svg"}
        skip_dirs = {".git", "__pycache__", ".ipynb_checkpoints", ".doctrees", "_build"}
        if ext in skip_exts:
            continue
        if any(d in fpath.parts for d in skip_dirs):
            continue
        files_by_ext[ext].append(fpath)

    category_map = {
        "courseware": {".ipynb", ".md", ".rst", ".html", ".htm"},
        "ppt": {".pptx", ".ppt"},
        "code": {".py"},
        "textbook": {".pdf", ".tex", ".epub"},
        "dataset": {".csv"},
        "reference": {".json", ".yaml", ".yml", ".txt", ".toml", ".jsonl"},
    }

    result = {}
    for cat, exts in category_map.items():
        cat_files = {}
        for ext in exts:
            if ext in files_by_ext:
                cat_files[ext] = files_by_ext[ext]
        result[cat] = cat_files

    return result


if __name__ == "__main__":
    print("=" * 60)
    print("EduMatrix FAISS Ingestion Pipeline")
    print("=" * 60)

    classified = classify_files()

    for cat_name in ["courseware", "ppt", "code", "textbook", "dataset", "reference"]:
        files_by_type = classified.get(cat_name, {})
        total_file_count = sum(len(v) for v in files_by_type.values())
        if total_file_count == 0:
            print(f"\n  Skipping {cat_name}: no files found")
            continue
        # Print what we found
        for ext, paths in files_by_type.items():
            print(f"  {ext}: {len(paths)} files")
        build_category_index(cat_name, CATEGORIES[cat_name], files_by_type)

    print(f"\n{'=' * 60}")
    print("INGESTION SUMMARY")
    print(f"{'=' * 60}")
    for cat, stats in sorted(TOTAL_STATS.items()):
        print(f"  {cat:20s}: {stats['files']:4d} files -> {stats['chunks']:6d} chunks")
    total_all = sum(s["chunks"] for s in TOTAL_STATS.values())
    total_files = sum(s["files"] for s in TOTAL_STATS.values())
    print(f"  {'TOTAL':20s}: {total_files:4d} files -> {total_all:6d} chunks")
    print(f"\nIndexes saved to {FAISS_DIR}")
